import json
import os
import tempfile
from lxml import etree
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.xmlchemy import OxmlElement
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

def link_to_slide(run, target_slide):
    
    r_id = run.part.relate_to(
        target_slide.part,
        RT.SLIDE,
    )
    
    rPr = run._r.get_or_add_rPr()
    
    hlinkClick = rPr.add_hlinkClick(r_id)
    hlinkClick.set('action', 'ppaction://hlinksldjump')

def unbullet(p):
    p._pPr.insert(
        0,
        etree.Element(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
        ),
    )
    p._element.get_or_add_pPr().set("marL", "0")
    p._element.get_or_add_pPr().set("indent", "0")

def titlify(p):
    """
    주어진 paragraph(p)에 대해 major theme fonts를 명시적으로 설정.
    - Latin: +mj-lt
    - East Asian: +mj-ea
    """
    # <a:defRPr> 요소 생성 또는 가져오기
    pPr = p._element.get_or_add_pPr()
    defRPr = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr")
    if defRPr is not None:
        # 기존 거 있으면 제거 (덮어쓰기 위해)
        pPr.remove(defRPr)

    defRPr = OxmlElement("a:defRPr")

    # <a:latin typeface="+mj-lt"/>
    latin = OxmlElement("a:latin")
    latin.set("typeface", "+mj-lt")

    # <a:ea typeface="+mj-ea"/>
    ea = OxmlElement("a:ea")
    ea.set("typeface", "+mj-ea")

    defRPr.append(latin)
    defRPr.append(ea)

    # defRPr 추가
    pPr.append(defRPr)


def orderify(p):
    """
    p.level 값을 기준으로 번호 스타일 설정
    """
    level = p.level

    # 원하는 스타일 매핑
    # 참고: https://learn.microsoft.com/en-us/dotnet/api/documentformat.openxml.drawing.autonumberschemevalues
    style_map = {
        0: "arabicPeriod",   # 1.
        1: "arabicParenR",   # 1)
        2: "alphaLcParenR",  # a)
        3: "alphaUcParenR",  # A)
        4: "romanLcParenR",  # i)
    }

    auto_num_type = style_map.get(level, "arabicPeriod")

    pPr = p._element.get_or_add_pPr()

    # 기존 불릿 제거
    for tag in ["a:buChar", "a:buAutoNum"]:
        el = pPr.find(f".//{tag}", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        if el is not None:
            pPr.remove(el)

    # buAutoNum 추가
    buAutoNum = OxmlElement("a:buAutoNum")
    buAutoNum.set("type", auto_num_type)
    pPr.append(buAutoNum)


def set_highlight(run, color):
    # get run properties
    rPr = run._r.get_or_add_rPr()
    # Create highlight element
    hl = OxmlElement("a:highlight")
    # Create specify RGB Colour element with color specified
    srgbClr = OxmlElement("a:srgbClr")
    setattr(srgbClr, "val", color)
    # Add colour specification to highlight element
    hl.append(srgbClr)
    # Add highlight element to run properties
    setattr(rPr, "lang", MSO_LANGUAGE_ID.ENGLISH_US)
    setattr(rPr, "altLang", MSO_LANGUAGE_ID.KOREAN)
    # lang="en-US" altLang="ko-KR"
    rPr.append(hl)
    latin = OxmlElement("a:latin")
    # <a:latin typeface="Consolas" panose="020B0609020204030204" pitchFamily="49" charset="0"/>
    setattr(latin, "typeface", "Consolas")
    setattr(latin, "charset", "0")
    rPr.append(latin)
    return run


def dict_shape(shape, placeholder=None):
    """
    주어진 shape 객체의 속성을 딕셔너리 형태로 반환합니다.
    """
    try:
        from_pl = json.loads(placeholder.name)
    except:
        from_pl = {}
    return {
        "name": shape.name or "",
        "top": shape.top or 0,
        "left": shape.left or 0,
        "width": shape.width or 0,
        "height": shape.height or 0,
        "right": shape.left + shape.width or 0,
        "bottom": shape.top + shape.height or 0,
        **from_pl,
    }


def clear_slides(prs):
    """
    while 루프를 사용하여 프레젠테이션의 모든 슬라이드를 삭제합니다.
    각 슬라이드에 대해 rId 관계를 삭제한 후, 슬라이드 ID 요소를 제거합니다.
    마지막에 임시 파일로 저장 후 재로드하여 내부 구조를 정리합니다.
    """
    # _sldIdLst는 슬라이드 ID들의 리스트입니다.
    while len(prs.slides._sldIdLst) > 0:
        slide_id = prs.slides._sldIdLst[0]
        # 슬라이드의 관계(rId)를 삭제합니다.
        prs.part.drop_rel(slide_id.rId)
        # 첫 번째 슬라이드를 삭제합니다.
        prs.slides._sldIdLst.remove(slide_id)
    # 내부 구조 정리를 위해 임시 파일에 저장 후 재로드
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
        temp_filename = tmp_file.name
    prs.save(temp_filename)
    new_prs = Presentation(temp_filename)
    os.remove(temp_filename)
    return new_prs
def boldify(run, width=12700, theme_color="accent3", alpha=0):
    rPr = run._r.get_or_add_rPr()

    # 기존 <a:ln> 제거
    for child in rPr.findall("./a:ln", namespaces=rPr.nsmap):
        rPr.remove(child)

    ln = OxmlElement("a:ln")
    ln.set("w", str(width))

    solidFill = OxmlElement("a:solidFill")
    schemeClr = OxmlElement("a:schemeClr")
    schemeClr.set("val", theme_color)

    # alpha 설정
    alpha_elem = OxmlElement("a:alpha")
    alpha_elem.set("val", str(alpha))
    schemeClr.append(alpha_elem)

    solidFill.append(schemeClr)
    ln.append(solidFill)

    rPr.insert(0, ln)
