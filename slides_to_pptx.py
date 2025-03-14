#!/usr/bin/env python3
"""
Slides to PowerPoint Converter

This script converts the structured JSON slide format into a PowerPoint presentation
using the python-pptx library.

Usage:
    # Convert JSON slides to PowerPoint
    python slides_to_pptx.py --file slides.json --export presentation.pptx
    
Requirements:
    - python-pptx: PowerPoint file generation library
    - Install with: pip install python-pptx
"""
import argparse
import json
import sys
import os
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def hex_to_rgb(hex_color: str) -> tuple:
    """
    Convert hex color to RGB tuple.
    
    Args:
        hex_color: Hex color code (e.g., "#FF0000")
        
    Returns:
        Tuple of (r, g, b) values
    """
    # Remove the "#" if present
    hex_color = hex_color.lstrip("#")
    # Convert to RGB
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def apply_text_style(text_frame, spans: List[Dict[str, Any]]) -> None:
    """
    Apply text styles to a text frame based on span information.
    
    Args:
        text_frame: The PowerPoint text frame to style
        spans: List of span dictionaries with text and style information
    """
    if not spans:
        return
    
    # Clear any existing text
    if len(text_frame.paragraphs) == 0:
        p = text_frame.add_paragraph()
    else:
        p = text_frame.paragraphs[0]
        p.clear()
    
    # Add spans with styling
    for span in spans:
        text = span.get("text", "")
        styles = span.get("styles", {})
        
        # Skip empty spans
        if not text:
            continue
        
        # Add a run for this span
        run = p.add_run()
        run.text = text
        
        # Apply styles
        if styles.get("bold"):
            run.font.bold = True
        
        if styles.get("italic"):
            run.font.italic = True
        
        if styles.get("underline"):
            run.font.underline = True
        
        if "color" in styles:
            # Convert hex color to RGB
            try:
                r, g, b = hex_to_rgb(styles["color"])
                run.font.color.rgb = RGBColor(r, g, b)
            except ValueError:
                print(f"Warning: Invalid color format: {styles['color']}", file=sys.stderr)
        
        if "link" in styles:
            run.hyperlink.address = styles["link"]


def add_paragraph_element(slide, element: Dict[str, Any], left: float, top: float, 
                          width: float, height: Optional[float] = None) -> float:
    """
    Add a paragraph element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The paragraph element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches (optional)
        
    Returns:
        The bottom position of the added element
    """
    # Create a text box
    txbox = slide.shapes.add_textbox(
        left=Inches(left), 
        top=Inches(top), 
        width=Inches(width), 
        height=Inches(height or 1.0)
    )
    
    # Apply text content and styling
    apply_text_style(txbox.text_frame, element.get("spans", []))
    
    # Word wrap
    txbox.text_frame.word_wrap = True
    
    # Estimate the height (this is approximate)
    # In a real implementation, you might need more sophisticated height calculation
    line_count = sum(len(span.get("text", "")) for span in element.get("spans", [])) // 50 + 1
    estimated_height = line_count * 0.3  # Rough estimate
    
    return top + estimated_height


def add_heading_element(slide, element: Dict[str, Any], left: float, top: float, 
                        width: float) -> float:
    """
    Add a heading element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The heading element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    # Create a text box for the heading
    txbox = slide.shapes.add_textbox(
        left=Inches(left), 
        top=Inches(top), 
        width=Inches(width), 
        height=Inches(1.0)
    )
    
    # Apply text content and styling
    level = element.get("level", 1)
    text_frame = txbox.text_frame
    
    # Apply text content and styling
    apply_text_style(text_frame, element.get("spans", []))
    
    # Style based on heading level
    p = text_frame.paragraphs[0]
    
    # Set font size based on heading level
    if level == 1:
        p.font.size = Pt(32)
    elif level == 3:
        p.font.size = Pt(24)
    elif level == 4:
        p.font.size = Pt(20)
    elif level == 5:
        p.font.size = Pt(18)
    elif level == 6:
        p.font.size = Pt(16)
    else:
        p.font.size = Pt(28)  # Default for other levels
    
    # Make headings bold
    p.font.bold = True
    
    # Estimate the height (this is approximate)
    estimated_height = 0.4  # Base height for headings
    
    return top + estimated_height


def add_list_element(slide, element: Dict[str, Any], left: float, top: float, 
                     width: float) -> float:
    """
    Add a list element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The list element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    # Create a text box for the list
    txbox = slide.shapes.add_textbox(
        left=Inches(left), 
        top=Inches(top), 
        width=Inches(width), 
        height=Inches(1.0)
    )
    
    text_frame = txbox.text_frame
    text_frame.word_wrap = True
    
    # Get list items
    items = element.get("items", [])
    
    # Clear default paragraph if it exists
    if text_frame.paragraphs:
        text_frame.clear()
    
    # Process list items
    for i, item in enumerate(items):
        # Add a paragraph for this item
        p = text_frame.add_paragraph()
        
        # Set indentation based on level
        level = item.get("level", 0)
        p.level = level
        
        # Add bullet character
        if element.get("style") == "number":
            p.text = f"{i+1}. {item.get('text', '')}"
        else:
            p.text = item.get("text", "")
            p.bullet = True
        
        # Process subitems if present
        subitems = item.get("subitems", [])
        for j, subitem in enumerate(subitems):
            sub_p = text_frame.add_paragraph()
            sub_p.level = subitem.get("level", level + 1)
            
            if element.get("style") == "number":
                sub_p.text = f"{i+1}.{j+1}. {subitem.get('text', '')}"
            else:
                sub_p.text = subitem.get("text", "")
                sub_p.bullet = True
    
    # Estimate the height (this is approximate)
    item_count = len(items) + sum(len(item.get("subitems", [])) for item in items)
    estimated_height = item_count * 0.3  # Rough estimate
    
    return top + estimated_height


def add_code_element(slide, element: Dict[str, Any], left: float, top: float, 
                     width: float) -> float:
    """
    Add a code block element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The code element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    # Create a text box for the code
    txbox = slide.shapes.add_textbox(
        left=Inches(left), 
        top=Inches(top), 
        width=Inches(width), 
        height=Inches(1.0)
    )
    
    # Get code text and language
    code_text = element.get("text", "")
    language = element.get("language", "")
    
    # Add language label if provided
    if language:
        code_text = f"{language}:\n{code_text}"
    
    # Set text
    text_frame = txbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = code_text
    
    # Style as code
    p.font.name = "Courier New"
    p.font.size = Pt(12)
    
    # Apply background color if specified
    if "background" in element.get("styles", {}):
        # Note: Setting shape fill color is more complex in python-pptx
        # This is a simplified version
        try:
            r, g, b = hex_to_rgb(element["styles"]["background"])
            txbox.fill.solid()
            txbox.fill.fore_color.rgb = RGBColor(r, g, b)
        except (ValueError, AttributeError):
            print(f"Warning: Couldn't set code block background color", file=sys.stderr)
    
    # Estimate the height (this is approximate)
    line_count = code_text.count('\n') + 1
    estimated_height = line_count * 0.25  # Rough estimate for code
    
    return top + estimated_height


def add_image_element(slide, element: Dict[str, Any], left: float, top: float, 
                      width: float) -> float:
    """
    Add an image element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The image element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    # Get image path
    image_path = element.get("path", "")
    
    if not image_path or not os.path.exists(image_path):
        print(f"Warning: Image not found: {image_path}", file=sys.stderr)
        return top + 0.5  # Return small height if image not found
    
    # Add the image
    try:
        img = slide.shapes.add_picture(
            image_path,
            left=Inches(left),
            top=Inches(top),
            width=Inches(width)
        )
        
        # Calculate height (maintaining aspect ratio)
        img_width = img.width
        img_height = img.height
        aspect_ratio = img_height / img_width
        height_in_inches = width * aspect_ratio / 914400  # Convert from EMUs to inches
        
        # Add caption if present
        caption = element.get("caption", "")
        if caption:
            caption_top = top + height_in_inches + 0.1
            txbox = slide.shapes.add_textbox(
                left=Inches(left), 
                top=Inches(caption_top), 
                width=Inches(width), 
                height=Inches(0.5)
            )
            txbox.text_frame.paragraphs[0].text = caption
            txbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            return caption_top + 0.5
        
        return top + height_in_inches
        
    except Exception as e:
        print(f"Error adding image: {e}", file=sys.stderr)
        return top + 0.5  # Return small height on error


def add_blockquote_element(slide, element: Dict[str, Any], left: float, top: float, 
                           width: float) -> float:
    """
    Add a blockquote element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The blockquote element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    # Create a text box with some indentation
    txbox = slide.shapes.add_textbox(
        left=Inches(left + 0.5),  # Add indentation
        top=Inches(top), 
        width=Inches(width - 1.0),  # Reduce width to account for indentation
        height=Inches(1.0)
    )
    
    # Apply text content and styling
    apply_text_style(txbox.text_frame, element.get("spans", []))
    
    # Style the blockquote with italic
    p = txbox.text_frame.paragraphs[0]
    p.font.italic = True
    
    # Apply blockquote color if specified
    if "color" in element.get("styles", {}):
        try:
            r, g, b = hex_to_rgb(element["styles"]["color"])
            p.font.color.rgb = RGBColor(r, g, b)
        except ValueError:
            print(f"Warning: Invalid color format for blockquote", file=sys.stderr)
    
    # Estimate the height (this is approximate)
    text_length = sum(len(span.get("text", "")) for span in element.get("spans", []))
    line_count = text_length // 40 + 1  # Rough estimate considering the reduced width
    estimated_height = line_count * 0.3
    
    return top + estimated_height


def add_table_element(slide, element: Dict[str, Any], left: float, top: float, 
                      width: float) -> float:
    """
    Add a table element to a slide.
    
    Args:
        slide: The PowerPoint slide
        element: The table element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    headers = element.get("headers", [])
    rows = element.get("rows", [])
    
    if not headers and not rows:
        return top  # No table data
    
    # Determine table dimensions
    cols = len(headers) or (rows[0] if rows else 0)
    if cols == 0:
        return top  # No columns
    
    # Create table
    table_rows = 1 + len(rows)  # Headers + data rows
    table = slide.shapes.add_table(
        rows=table_rows,
        cols=cols,
        left=Inches(left),
        top=Inches(top),
        width=Inches(width),
        height=Inches(0.4 * table_rows)  # Approximate height
    )
    
    # Add headers
    for i, header in enumerate(headers):
        if i < cols:
            cell = table.table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx < cols:
                cell = table.table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
    
    # Return the position below the table
    return top + (0.4 * table_rows)


def add_content_element(slide, element: Dict[str, Any], left: float, top: float, 
                        width: float) -> float:
    """
    Add a content element to a slide based on its type.
    
    Args:
        slide: The PowerPoint slide
        element: The content element dictionary
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        
    Returns:
        The bottom position of the added element
    """
    element_type = element.get("type", "")
    
    if element_type == "paragraph":
        return add_paragraph_element(slide, element, left, top, width)
    
    elif element_type == "heading":
        return add_heading_element(slide, element, left, top, width)
    
    elif element_type == "list":
        return add_list_element(slide, element, left, top, width)
    
    elif element_type == "code":
        return add_code_element(slide, element, left, top, width)
    
    elif element_type == "image":
        return add_image_element(slide, element, left, top, width)
    
    elif element_type == "blockquote":
        return add_blockquote_element(slide, element, left, top, width)
    
    elif element_type == "table":
        return add_table_element(slide, element, left, top, width)
    
    # Unknown element type
    print(f"Warning: Unknown element type: {element_type}", file=sys.stderr)
    return top


def create_slide(prs, slide_data: Dict[str, Any]) -> None:
    """
    Create a PowerPoint slide from slide data.
    
    Args:
        prs: The PowerPoint presentation object
        slide_data: The slide data dictionary
    """
    # Get slide properties
    title = slide_data.get("title", "")
    layout = slide_data.get("layout", "title_and_content")
    slide_id = slide_data.get("id", "")
    
    # Choose PowerPoint layout based on our layout type
    if layout == "title_only":
        slide_layout = prs.slide_layouts[5]  # Title Only
    elif layout == "two_content":
        slide_layout = prs.slide_layouts[3]  # Two Content
    else:
        slide_layout = prs.slide_layouts[1]  # Title and Content
    
    # Add a slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title if present
    if title and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
        slide.shapes.title.text = title
    
    # Process content based on layout
    if layout == "two_content" and "columns" in slide_data:
        # For two-column layout
        columns = slide_data.get("columns", [])
        
        if len(columns) >= 2:
            # Calculate column positions
            col1_width = columns[0].get("width", 0.5)
            col2_width = columns[1].get("width", 0.5)
            
            slide_width = 10.0  # Standard slide width in inches
            margin = 0.5  # Margin in inches
            
            col1_left = margin
            col1_width_inches = (slide_width - 2 * margin) * col1_width
            
            col2_left = col1_left + col1_width_inches + margin/2
            col2_width_inches = (slide_width - 2 * margin) * col2_width
            
            # Process first column content
            top = 1.5  # Start below title
            for element in columns[0].get("content", []):
                top = add_content_element(slide, element, col1_left, top, col1_width_inches) + 0.2
            
            # Process second column content
            top = 1.5  # Reset top position for second column
            for element in columns[1].get("content", []):
                top = add_content_element(slide, element, col2_left, top, col2_width_inches) + 0.2
    else:
        # For single-column layout
        content = slide_data.get("content", [])
        margin = 0.5  # Margin in inches
        slide_width = 10.0  # Standard slide width in inches
        content_width = slide_width - 2 * margin
        
        # Process content elements
        top = 1.5  # Start below title
        for element in content:
            top = add_content_element(slide, element, margin, top, content_width) + 0.2


def create_presentation(slides_data: Dict[str, Any]) -> Presentation:
    """
    Create a PowerPoint presentation from structured slide data.
    
    Args:
        slides_data: Dictionary containing presentation data
        
    Returns:
        A PowerPoint presentation object
    """
    # Create a blank presentation
    prs = Presentation()
    
    # Get presentation metadata
    metadata = slides_data.get("presentation", {}).get("metadata", {})
    title = metadata.get("title", "Untitled Presentation")
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = metadata.get("author", "")
    
    # Create slides
    slides = slides_data.get("presentation", {}).get("slides", [])
    for slide_data in slides:
        create_slide(prs, slide_data)
    
    return prs


def main():
    """
    Main function to process command line arguments and create the PowerPoint.
    """
    parser = argparse.ArgumentParser(description='Convert structured slide JSON to PowerPoint presentation')
    
    parser.add_argument('--file', '-f', required=True, help='Path to the slide JSON file')
    parser.add_argument('--export', '-e', required=True, help='Path to export the PowerPoint file')
    
    args = parser.parse_args()
    
    # Read the JSON file
    try:
        with open(args.file, 'r', encoding='utf-8') as f:
            slides_data = json.load(f)
    except Exception as e:
        print(f"Error reading JSON file: {e}", file=sys.stderr)
        sys.exit(1)
    
    # Create the presentation
    try:
        prs = create_presentation(slides_data)
        
        # Save the presentation
        prs.save(args.export)
        print(f"Successfully created PowerPoint presentation: {args.export}")
    except Exception as e:
        print(f"Error creating PowerPoint: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()