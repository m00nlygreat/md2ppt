from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "글머리 기호 제거 예제"

# 본문 placeholder
content = slide.placeholders[1]
tf = content.text_frame

# 첫 번째 단락 (기본적으로 bullet 있음)
tf.text = "이건 기본 bullet"

# 두 번째 단락 (bullet 제거)
p = tf.add_paragraph()
p.text = "이건 bullet 없음"
p.level = 0
p._element.get_or_add_pPr().set("marL", "0")  # 선택적으로 왼쪽 여백도 제거
p._element.get_or_add_pPr().set("indent", "0")  # 선택적으로 들여쓰기 제거
p._bullet = False  # 핵심: bullet 제거

# 저장
prs.save("no_bullet.pptx")
print("✅ bullet 제거된 PPT 생성 완료")
