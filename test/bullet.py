from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn


def remove_bullet_and_indent(paragraph):
    """
    주어진 Paragraph에서 bullet 속성과 들여쓰기를 제거합니다.
    bullet 관련 태그(a:buChar, a:buAutoNum, a:buBlip)가 있으면 제거하며,
    paragraph_format의 left_indent와 first_line_indent를 0으로 설정합니다.
    """
    # XML pPr 요소 가져오기 (없으면 추가됨)
    pPr = paragraph._p.get_or_add_pPr()
    # bullet 관련 속성을 찾아 제거
    for bullet_tag in ["a:buChar", "a:buAutoNum", "a:buBlip"]:
        bullet = pPr.find(qn(bullet_tag))
        if bullet is not None:
            pPr.remove(bullet)

    # 들여쓰기 제거: 왼쪽 들여쓰기와 첫 줄 들여쓰기 0으로 설정
    para_format = paragraph.paragraph_format
    para_format.left_indent = Pt(0)
    para_format.first_line_indent = Pt(0)


# 새로운 프레젠테이션 생성
prs = Presentation()

# 기본 프레젠테이션의 slide_layouts[1]은 일반적으로 "Title and Content" 레이아웃입니다.
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

# (선택) 제목 placeholder에 간단한 제목 입력
slide.shapes.title.text = "제목 없음"

# 콘텐츠 placeholder의 TextFrame 접근 및 기존 내용을 제거
content_placeholder = slide.placeholders[1]
text_frame = content_placeholder.text_frame
text_frame.clear()

# 새 paragraph 추가 후 임의의 텍스트 입력
p = text_frame.add_paragraph()
p.text = "첫 줄에 들어갈 텍스트입니다."

# bullet과 왼쪽 들여쓰기를 제거
remove_bullet_and_indent(p)

# 프레젠테이션 저장
prs.save("result.pptx")
