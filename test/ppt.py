import pptx

p = pptx.Presentation("./ref.pptx")
# obj = p.slides[2].shapes[0].text_frame.paragraphs[0].runs[0].font.size
obj = p.slide_masters[0].placeholders[0].text_frame.paragraphs[0].runs[0].font.size


def print_obj(obj):
    """
    객체의 속성 중에서 호출 가능한(callable) 메서드를 제외한 속성들을 key: value 형식으로 출력합니다.
    private 속성(이름이 '_'로 시작하는)은 생략합니다.
    """
    for key in dir(obj):
        # private 속성은 무시합니다.
        if key.startswith('_'):
            continue
        try:
            value = getattr(obj, key)
        except Exception:
            continue
        # callable(메서드 등)은 제외합니다.
        if callable(value):
            continue
        try:
            print(f"{key}: {value}")
        except Exception:
            print(f"{key}: <non-printable>")
            
print(type(obj))
print(dir(obj))
print(obj)
print_obj(obj)
