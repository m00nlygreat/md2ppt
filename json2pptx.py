import json, os, re, warnings, argparse
from enum import Enum
from pptx import Presentation
from PIL import Image
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from utils.util import unbullet, orderify, set_highlight, dict_shape, clear_slides, link_to_slide
from utils.expand import expand
from utils.code_highlight import highlight_code, process_codes

pholder = 0

TEXT_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

def get_slide_layout_enum(prs):

    layout_members = {}
    for idx, layout in enumerate(prs.slide_layouts):
        # layout.name 속성이 없을 경우에는 기본 이름 사용
        raw_name = getattr(layout, "name", f"LAYOUT_{idx}")
        # 유효한 enum 멤버 이름으로 변환 (대문자로, 숫자로 시작하는 경우 앞에 '_' 추가)
        member_name = re.sub(r"\W|^(?=\d)", "_", raw_name.upper())
        if not member_name:
            member_name = f"LAYOUT_{idx}"
        # 중복 방지를 위해 이미 존재하면 인덱스를 추가
        if member_name in layout_members:
            member_name = f"{member_name}_{idx}"
        layout_members[member_name] = idx
    SlideLayoutEnum = Enum("SlideLayoutEnum", layout_members)
    return SlideLayoutEnum

def convert_json_to_pptx(prs, data, layouts, toc=1):
    
    def add_toc_item(paragraph, item):
        title_run = paragraph.add_run()
        title_run.text = item.get("title", "") + '\t'
        index_run = paragraph.add_run()
        slide_no = item.get("index", 0)+1
        index_run.text = str(slide_no)
        link_to_slide(index_run, prs.slides[slide_no-1])
        
        # index_run.hyperlink.address = f'#slide=id.p{slide_no}'
    
    slides_data = data.get("slides", [])
    prev_title = None
    for current_slide_no, slide in enumerate(slides_data):
        layout_name_from_json = slide.get("layout", "title_and_content").upper()
        try:
            layout_index = layouts[layout_name_from_json].value
        except KeyError:
            # 레이아웃 이름이 Enum에 없을 경우 기본 레이아웃 사용
            print(f'Layout "{layout_name_from_json}" not found. Using default layout.')
            layout_index = None
            for layout in layouts:
                if layout.name == "TITLE_AND_CONTENT":
                    layout_index = layout.value
                    break

        slide_layout_idx = prs.slide_layouts[layout_index]
        # print(slide_layout_idx.name)
        current_slide = prs.slides.add_slide(slide_layout_idx)

        p_map = {}
        for i, pl in enumerate(current_slide.placeholders):
            p_map.update({i: slide_layout_idx.placeholders[i].placeholder_format.idx})

        # 제목을 설정합니다.
        title = slide.get("title", False)
        title_shape = current_slide.shapes.title
        p = title_shape.text_frame.paragraphs[0]
        if title:
            runs = title.get("runs", [])
            process_runs(runs, p)
            prev_title = runs
        else: 
            if isinstance(prev_title, list):
                process_runs(prev_title, p)            

        # Placeholder에 토큰을 처리합니다.
        # grow 룰을 적용하기 위한 타이틀 제외 shape (실제 추가된 순서로)를 모아둠
        
        shapes_no_title = []
        pl_after = False
        global pholder_no
        pholder_no = 0
        placeholder_count = len(current_slide.placeholders)
        for pholder_data in slide.get("placeholders", []):
            pholder_no += 1
            if placeholder_count >= pholder_no:
                try:
                    current_placeholder = current_slide.placeholders[p_map[pholder_no]]
                except:
                    pass

                for token in pholder_data:
                    pl_after = process_token(current_placeholder, token, current_slide)
                    # image이면 picture shape, 텍스트이면 placeholder가 들어있게 될 것.
            else:
                print(f"Error: Placeholder index {pholder_no} exceeds available placeholders.")
            if pl_after:
                shapes_no_title.append(pl_after)
        
        # 남는 공간을 차지하도록 채웁니다.
        # Title placeholder는 무적권 0번이어야 해
        # shapes_no_title = [sh for i, sh in enumerate(current_slide.shapes) if i != 0]
        
        shapes = []
        
        for i, shape in enumerate(shapes_no_title):
            if current_slide.shapes.title == shape:
                print(current_slide.shapes.title.text_frame.text)
                continue
            placeholder = current_slide.slide_layout.placeholders[i+1] # 추후 고쳐줘야 한다. 에러나서 안되기 때문에.
            shapes.append(dict_shape(shape, placeholder))
        
        # align 적용
        for i, shape in enumerate(shapes):
            align = shape.get("align", False)
            if align:
                width, height = shape['width'], shape['height']
                new_sizloc = calc_align(current_placeholder, width, height, align)
                foo = shapes_no_title[i]
                foo.left = new_sizloc['left']
                foo.top = new_sizloc['top']
                foo.width = new_sizloc['width']
                foo.height = new_sizloc['height']

        # Shape의 size는 한 번에 대입해줘야 한다. 하나씩 변경하면 0으로 초기화됨
        for i,shape in enumerate(shapes_no_title):
            grow = shapes[i].get("grow", False)
            if grow:
                foo_shp = shapes[i]
                l, r, a, b = expand(shapes, i, prs).values()
                if grow in [1,2,3]:
                    foo_shp['height'] += b
                if grow in [1,4,7]:
                    foo_shp['left'] -= l
                    foo_shp['width'] += l + r
                if grow in [3,6,9]:
                    foo_shp['width'] += r
                if grow in [7,8,9]:
                    foo_shp['top'] -= a
                    foo_shp['height'] += a + b
                if grow == 5:
                    foo_shp['left'] -= l
                    foo_shp['top'] -= a
                    foo_shp['width'] += l + r
                    foo_shp['height'] += a + b
                # 한꺼번에 지정해야한다. 지정할 때 나머지가 0으로 초기화되기 때문. left/top이 width/height와 세트인 것으로 보임.
                shape.left = foo_shp['left']
                shape.top = foo_shp['top']
                shape.width = foo_shp['width']
                shape.height = foo_shp['height']
    
    # TOC 슬라이드 추가
    if toc:
        try:
            layout_name = 'toc'.upper()
            layout_idx = layouts[layout_name].value
        except:
            # 레이아웃 이름이 Enum에 없을 경우 기본 레이아웃 사용
            print(f'Layout "TOC" not found. Using default layout.')
            layout_name = 'two_content'.upper()
            layout_idx = layouts[layout_name].value
        
        toc_slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        toc_slide.shapes.title.text = "Table of Contents"
        toc_placeholder = toc_slide.placeholders[1]
        toc_data = data.get("toc", []).get("chapters", False)
        if toc_data and toc>=1:
            for item in toc_data:
                p = define_paragraph(toc_placeholder)
                p.level = 0
                add_toc_item(p, item)
                modules = item.get("modules", False)
                if modules and toc >=2:
                    for item in modules:
                        p = define_paragraph(toc_placeholder)
                        p.level = 1
                        add_toc_item(p, item)
            
    
def calc_align(p, width, height, align=5):

    # align 값을 정수로 변환 시도, 실패하거나 1~9 범위가 아니면 기본값 5 사용
    
    try:
        align_val = int(align)
    except Exception:
        align_val = 5
    if align_val < 1 or align_val > 9:
        align_val = 5

    # numpad 정렬 값에 따른 가로 정렬 계수: 왼쪽=0, 가운데=0.5, 오른쪽=1
    if align_val in (1, 4, 7):
        factor_x = 0
    elif align_val in (2, 5, 8):
        factor_x = 0.5
    else:  # align_val in (3, 6, 9)
        factor_x = 1

    # numpad 정렬 값에 따른 세로 정렬 계수: 위=0, 가운데=0.5, 아래=1
    if align_val in (7, 8, 9):
        factor_y = 0
    elif align_val in (4, 5, 6):
        factor_y = 0.5
    else:  # align_val in (1, 2, 3)
        factor_y = 1

    # Placeholder의 좌표와 크기
    p_left = p.left
    p_top = p.top
    p_width = p.width
    p_height = p.height
    p_ratio = p_width / p_height

    # 이미지의 원본 크기 및 비율
    i_width, i_height = width, height
    i_ratio = i_width / i_height

    if i_ratio < p_ratio:
        # 이미지가 placeholder보다 상대적으로 좁은 경우: 높이를 맞추고 너비를 조절
        new_width = p_height * i_ratio
        # 남는 가로 공간을 factor_x에 따라 오프셋 적용
        new_left = p_left + (p_width - new_width) * factor_x
        align_to = {
            "left": int(new_left),
            "top": int(p_top),  # 세로는 꽉 채움
            "width": int(new_width),
            "height": int(p_height),
        }
    else:
        # 이미지가 placeholder보다 상대적으로 넓은 경우: 너비를 맞추고 높이를 조절
        new_height = p_width / i_ratio
        # 남는 세로 공간을 factor_y에 따라 오프셋 적용
        new_top = p_top + (p_height - new_height) * factor_y
        align_to = {
            "left": int(p_left),  # 가로는 꽉 채움
            "top": int(new_top),
            "width": int(p_width),
            "height": int(new_height),
        }
    return align_to

def process_token(current_placeholder, token, current_slide):

    match(token.get("type", "")):
        case "paragraph":
            p = define_paragraph(current_placeholder)
            unbullet(p)
            process_runs(token.get("runs", []), p)
        case "heading":
            p = define_paragraph(current_placeholder)
            # unbullet(p)
            # titlify(p)
            # p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
            # p.level = token.get("depth", 0)
            match(token.get("level", 3)):
                case 3:
                    p.level = 8
                case _:
                    p.level = 7
            process_runs(token.get("runs", []), p)
        case "block_quote":
            p = define_paragraph(current_placeholder)
            p.level = 6
            process_runs(token.get("runs", []), p)
        case "code":
            a = current_placeholder.text_frame.add_paragraph()
            a.level = 5
            unbullet(a)
            p = define_paragraph(current_placeholder)
            p.level = 5
            lang = token.get("lang", None)
            code = token.get("raw", False)
            highlighted = highlight_code(code, lang)
            process_codes(highlighted, p)
                # p.text = token.get("lang","plaintext")+'\n'+ token.get("raw", "")
        case "image":
            url = token.get("url", "")
            
            try:
                i = Image.open(url)
                            
                global pholder_no
                
                dynloc = {"order": pholder_no}

                try:
                    align_dict = json.loads(current_slide.slide_layout.placeholders[pholder_no].name)
                    dynloc.update(align_dict)
                except:
                    # print('Error: Placeholder name is not JSON format.')
                    pass
                
                align_to = calc_align(current_placeholder, i.width, i.height , dynloc.get("align",5))

                try:
                    current_placeholder.insert_picture(url)
                except Exception as e:

                    sp = current_placeholder._element
                    sp.getparent().remove(sp)

                    current_placeholder = current_slide.shapes.add_picture(url, **align_to)
                    
            except:
                print(f"Error: Image '{url}' not found or invalid.")

        case "list":
            children = token.get("children", [])
            if children:
                for child in children:
                    p = define_paragraph(current_placeholder)
                    # print(child.get("type", ""))
                    p.level = child.get("depth", 0)
                    process_runs(child.get("runs", []), p)
                    if child.get("ordered", False):
                        orderify(p)
        case "table":
            process_table(current_placeholder, token, current_slide)
            # current_placeholder = process_table(current_placeholder, token, current_slide)
            # align, grow에 포함시켜야 하기 때문에 이렇게 돼야 하지만, table 자체의 height를 1000 기본값으로 하기 때문에 제대로 작동 안함.
        case _ :
            print(token.get("type", ""))
            pass
    return current_placeholder

def process_table(current_placeholder, token, current_slide):
    import unicodedata

    def visual_length(s):
        return sum(2 if unicodedata.east_asian_width(c) in 'WF' else 1 for c in str(s))

    def get_column_weights(head_data, body_data):
        num_cols = len(head_data)
        weights = [0] * num_cols
        for i, cell in enumerate(head_data):
            text = "".join(run.get("text", "") for run in cell.get("runs", []))
            weights[i] = max(weights[i], visual_length(text))
        for row in body_data:
            for i, cell in enumerate(row):
                text = "".join(run.get("text", "") for run in cell.get("runs", []))
                weights[i] = max(weights[i], visual_length(text))
        return weights

    def dynamic_cap(num_cols, base=0.4):
        if num_cols <= 2:
            return 0.9
        elif num_cols == 3:
            return 0.6
        elif num_cols == 4:
            return 0.5
        return base

    def normalize_with_cap(weights, cap):
        total = sum(weights)
        raw_ratios = [w / total for w in weights] if total > 0 else [1/len(weights)] * len(weights)
        capped = [min(r, cap) for r in raw_ratios]
        remainder = 1.0 - sum(capped)
        flexible_indices = [i for i, r in enumerate(raw_ratios) if r <= cap]
        flexible_sum = sum(raw_ratios[i] for i in flexible_indices) or 1
        for i in flexible_indices:
            capped[i] += remainder * (raw_ratios[i] / flexible_sum)
        return capped

    def write_cell(cell_data, cell):
        p = define_paragraph(cell)
        process_runs(cell_data.get("runs", []), p)
        align = cell_data.get("align", "left")
        if align:
            p.alignment = TEXT_ALIGN[align]

    sizloc = {
        "left": current_placeholder.left,
        "top": current_placeholder.top,
        "width": current_placeholder.width,
        "height": 1000,
    }
    sp = current_placeholder._element
    sp.getparent().remove(sp)

    head_data = token.get("head", [])
    body_data = token.get("body", [])

    rows_count = len(body_data) + 1
    cols_count = len(head_data)

    shape = current_slide.shapes.add_table(
        rows_count,
        cols_count,
        **sizloc
    )

    # 테이블의 스타일을 설정
    shape._element.graphic.graphicData.tbl[0][-1].text = '{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}'

    table = shape.table

    for index, cell_data in enumerate(head_data):
        cell = table.cell(0, index)
        write_cell(cell_data, cell)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for row_index, row in enumerate(body_data):
        for col_index, cell_data in enumerate(row):
            cell = table.cell(row_index+1, col_index)
            write_cell(cell_data, cell)
            cell.vertical_anchor = MSO_ANCHOR.TOP

    # 열 너비 자동 계산 적용
    weights = get_column_weights(head_data, body_data)
    ratios = normalize_with_cap(weights, cap=dynamic_cap(len(weights)))
    total_width = sizloc["width"]
    for i, ratio in enumerate(ratios):
        table.columns[i].width = int(total_width * ratio)
    
    return shape

def define_paragraph(placeholder):
    """
    Placeholder에서 첫 번째 단락을 가져오고, 텍스트가 비어있으면 새 단락을 추가합니다.
    """
    if placeholder.text_frame.paragraphs[0].text != "":
        paragraph = placeholder.text_frame.add_paragraph()
    else:
        paragraph = placeholder.text_frame.paragraphs[0]
    return paragraph

def process_runs(runs, paragraph):
    """
    주어진 runs 리스트를 사용하여 paragraph에 텍스트와 스타일을 설정합니다.
    각 run은 텍스트와 스타일 정보를 포함하는 딕셔너리입니다.
    """
    for run in runs:
        r = paragraph.add_run()
        r.text = run.get("text", "")
        font = r.font
        if 'bold' in run:
            font.bold = True
            font.color.theme_color = MSO_THEME_COLOR.ACCENT_3
        if 'italic' in run:
            font.italic = True
            font.underline = True
        if 'monospace' in run:
            r = set_highlight(r, 'EEEEEE')
            r.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
            # r.font.color.rgb = RGBColor(248, 104, 107)
            # print(font.size)
            # 현재 폰트 사이즈를 알아내는 게 쉽지 않다.
        if 'hyperlink' in run:
            r.hyperlink.address = run.get("hyperlink", "https://google.com")

def add_title_slide(prs, frontmatter):
    """
    주어진 Presentation 객체(prs)에 제목 슬라이드를 추가합니다.
    frontmatter는 제목 슬라이드에 표시할 정보를 포함하는 딕셔너리입니다.
    """
    title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃
    slide = prs.slides.add_slide(title_slide_layout)
    pp = prs.core_properties

    # 이후 placeholder 이름과 md-frontmatter key를 매치할 것
    title = frontmatter.get("title", False)
    subtitle = frontmatter.get("subtitle", False)
    author = frontmatter.get("author", False)   
    
    slide.shapes.title.text = title if title else "제목없음"
    pp.title = title if title else "Powerpoint 프레젠테이션"
    
    if subtitle:
        first_slide_subtitle = slide.placeholders[1]
        first_slide_subtitle = subtitle
        pp.subtitle = subtitle
    if author:
        pp.author = author

def main(data=None):
    parser = argparse.ArgumentParser(
        description="Convert JSON to PPTX using python-pptx"
    )
    parser.add_argument(
        "-r",
        "--ref",
        type=str,
        default=None,
        help="Reference PPTX file path (if provided and non-empty, slides will be cleared)",
    )
    parser.add_argument(
        "-o",
        "--output",
        type=str,
        default="output.pptx",
        help="Output PPTX file name (default: output.pptx)",
    )
    parser.add_argument(
        "-i", "--input", type=str, default=None, help="Input JSON file path"
    )
    parser.add_argument(
        "--return-pptx", action="store_true", help="Return Presentation object instead of saving to file"
    )
    args = parser.parse_args()

    # 환경 변수에서 매개변수 가져오기
    ref_from_env = os.environ.get("JSON2PPTX_REF", "")
    output_from_env = os.environ.get("JSON2PPTX_OUTPUT", "")
    return_pptx_from_env = os.environ.get("JSON2PPTX_RETURN_PPTX", "")

    # 환경 변수에서 가져온 값으로 args 업데이트
    if ref_from_env and not args.ref:
        args.ref = ref_from_env
    if output_from_env and not args.output:
        args.output = output_from_env
    if return_pptx_from_env and not args.return_pptx:
        args.return_pptx = True

    # JSON 데이터 로딩: 딕셔너리를 직접 전달받은 경우 우선 사용
    if data is None:
        if args.input is None:
            print(
                "Error: Please provide input JSON file path with -i/--input option or pass a dictionary to main()."
            )
            return
        else:
            if not os.path.exists(args.input):
                print(f"Error: JSON file '{args.input}' does not exist.")
                return
            with open(args.input, "r", encoding="utf-8") as f:
                data = json.load(f)

    # 참조 PPTX 파일이 지정되었고 존재하면 이를 사용합니다.
    if args.ref and os.path.exists(args.ref):
        prs = Presentation(args.ref)
        # 슬라이드가 있다면 완전히 제거 (첫 슬라이드도 포함)
        if len(prs.slides) > 0:
            prs = clear_slides(prs)
    else:
        # 참조 파일이 없으면 새 프레젠테이션 생성
        prs = Presentation()
        
    add_title_slide(prs, data['frontmatter'])

    layouts = get_slide_layout_enum(prs)
    # for layout in layouts:
        # print(f"{layout.name} = {layout.value}")

    # JSON 데이터를 기반으로 PPTX 변환 로직 실행
    convert_json_to_pptx(prs, data, layouts=layouts)

    # Presentation 객체 반환 모드
    if args.return_pptx:
        return prs

    # 출력 PPTX 파일 저장
    prs.save(args.output)
    print(f"PPTX file saved as {args.output}")

if __name__ == "__main__":
    main()
